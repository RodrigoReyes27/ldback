from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
import os  # Import the os module to work with file paths
import requests

powerpoint_file = requests.get("https://firebasestorage.googleapis.com/v0/b/frida-research.appspot.com/o/uploaded_files%2F326e5185-9290-4810-a7d2-fd1fc87b8d68.pptx?alt=media&token=dbf15c40-5309-4a0d-b022-abc759524870")

prs = Presentation(BytesIO(powerpoint_file.content))

n = 0

# Specify the directory where you want to save the images
output_directory = "./functions/app/documentos/assets/"

def write_image(shape):
    global n
    image = shape.image
    # ---get image "file" contents---
    image_bytes = image.blob
    # ---make up a name for the file, e.g. 'image.jpg'---
    image_filename = 'image{:03d}.{}'.format(n, image.ext)
    # Specify the full path to save the image
    image_path = os.path.join(output_directory, image_filename)
    n += 1
    print(image_path)
    with open(image_path, 'wb') as f:
        f.write(image_bytes)

def visitor(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for s in shape.shapes:
            visitor(s)
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        write_image(shape)

def iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            visitor(shape)

iter_picture_shapes(prs)

# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_runs = ""

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs += run.text
            text_runs += "\n"

print(text_runs)
from flask import jsonify, request
from pptx import Presentation
# from google.cloud import storage
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from firebase_admin import storage

# from os import getenv
# from dotenv import load_dotenv

import os  # Import the os module to work with file paths

from . import documentos_blueprint
from .. import FIREBASE_CONFIG

# load_dotenv()

@documentos_blueprint.route("/parse_pptx", methods=["POST"])
def parsear_pptx_handle():
    bucket_name = FIREBASE_CONFIG['STORAGE_BUCKET']
    storage_client = storage.Client(project=FIREBASE_CONFIG['PROJECT_ID'])
    bucket = storage_client.bucket(bucket_name)
    blob_name = request.args.get('file_name')
    blob = bucket.blob(blob_name)
    contents = blob.download_as_bytes()
    prs = Presentation(BytesIO(contents))

    text_runs = extract_text(prs)
    imglist = extract_images(prs, bucket_name)

    return jsonify(msg=text_runs, imglist=imglist)

def extract_text(prs):
    text_runs = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs += run.text
                text_runs += "\n"
    return text_runs

def extract_images(prs, bucket_name):
    imglist = []
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for s in shape.shapes:
                    if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        imglist.append(upload_image(s, n, bucket_name))
                        n += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                imglist.append(upload_image(shape, n, bucket_name))
                n += 1
    return imglist

def upload_image(shape, n, bucket_name):
    image = shape.image
    image_bytes = image.blob
    image_filename = 'jairexampleimage{:03d}.{}'.format(n, image.ext)
    storage_client = storage.Client(project=FIREBASE_CONFIG['PROJECT_ID'])
    bucket = storage_client.bucket(bucket_name)
    blob = bucket.blob(image_filename)
    blob.upload_from_string(image_bytes)
    return image_filename

# @user_blueprint.route("/parse_pptx", methods=["POST"])
# def parsear_pptx_handle():

#     """Downloads a blob into memory."""
#     # The ID of your GCS bucket
#     bucket_name = getenv("BUCKET_NAME")

#     # The ID of your GCS object
#     storage_client = storage.Client(project = getenv("GOOGLE_PROJECT_ID"))

#     bucket = storage_client.bucket(bucket_name)

#     blob_name = request.args.get('file_name')

#     blob = bucket.blob(blob_name)
#     contents = blob.download_as_bytes()
#     prs = Presentation(BytesIO(contents))

#     n = 0
#     finimglist = []

#     # Specify the directory where you want to save the images

#     finimglist = iter_picture_shapes(prs)

#     # text_runs will be populated with a list of strings,
#     # one for each text run in presentation
#     text_runs = ""

#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if not shape.has_text_frame:
#                 continue
#             for paragraph in shape.text_frame.paragraphs:
#                 for run in paragraph.runs:
#                     text_runs += run.text
#                 text_runs += "\n"


#     return jsonify(msg=text_runs)

# def write_image(shape, imglist):
#         global n
#         image = shape.image
#         # ---get image "file" contents---
#         image_bytes = image.blob
#         # ---make up a name for the file, e.g. 'image.jpg'---
#         image_filename = 'image{:03d}.{}'.format(n, image.ext)
#         # Specify the full path to save the image
#         image_path = os.path.join(output_directory, image_filename)
#         n += 1
#         storage_client = storage.Client()
#         bucket = storage_client.bucket(getenv("BUCKET_NAME"))
#         blob = bucket.blob("{image_filename}")

#         imglist.append(image_filename)
        
#         with open(image_path, 'wb') as f:
#             f.write(image_bytes)
#             blob.upload_from_string(f)


# def visitor(shape, imglist):
#     if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
#         for s in shape.shapes:
#             visitor(s, imglist)
#     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#         write_image(shape, imglist)

# def iter_picture_shapes(prs):
#     imglist = []
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             visitor(shape, imglist)
#             return imglist
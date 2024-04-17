from application.parser import IParser, ParsingResult
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO

class PPTXParser(IParser): 
    ...
    def parse(self, file: BytesIO) -> ParsingResult:
        prs = Presentation(file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    if paragraph.text == "":
                        continue
                    text_runs.append(paragraph.text)

        # imglist = []
        # n = 0
        # for slide in prs.slides:
        #     for shape in slide.shapes:
        #         if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        #             for s in shape.shapes:
        #                 if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
        #                     imglist.append()
        #                     n += 1
        #         elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        #             imglist.append()
        #             n += 1

        return ParsingResult(text=text_runs,images=None)